#!/usr/bin/env python3
"""
Convert markdown file to PowerPoint presentation
Usage: python md_to_pptx.py <input.md>
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
import sys
import os

def parse_markdown_slides(md_file):
    """Parse markdown file and extract slides"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide separator
    slides = content.split('---')

    # Remove the frontmatter (first section with marp config)
    slides = [s.strip() for s in slides if s.strip()]
    if slides and 'marp:' in slides[0]:
        slides = slides[1:]

    parsed_slides = []
    for slide in slides:
        if not slide.strip():
            continue

        lines = slide.strip().split('\n')

        # Extract title (lines starting with # or ##)
        title = ""
        content_lines = []

        for line in lines:
            if line.startswith('# '):
                title = line.lstrip('#').strip()
            elif line.startswith('## '):
                title = line.lstrip('#').strip()
            elif line.startswith('### '):
                # Subtitle/section header
                content_lines.append(line.lstrip('#').strip())
            else:
                content_lines.append(line)

        # Join content and clean up
        content = '\n'.join(content_lines).strip()

        parsed_slides.append({
            'title': title,
            'content': content
        })

    return parsed_slides

def create_presentation(slides, output_file):
    """Create PowerPoint presentation from parsed slides"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        # Use title and content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Set content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        # Parse content and add formatted text
        content = slide_data['content']

        # Split content into lines
        lines = content.split('\n')

        for i, line in enumerate(lines):
            if not line.strip():
                continue

            # Add paragraph
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            # Handle bullet points
            if line.strip().startswith('- '):
                text = line.strip()[2:]
                p.text = text
                p.level = 0
                p.font.size = Pt(18)
            elif line.strip().startswith('  - '):
                text = line.strip()[2:]
                p.text = text
                p.level = 1
                p.font.size = Pt(16)
            # Handle numbered lists
            elif re.match(r'^\d+\.', line.strip()):
                text = re.sub(r'^\d+\.\s*', '', line.strip())
                p.text = text
                p.level = 0
                p.font.size = Pt(18)
            # Handle bold text (markdown **text**)
            elif line.strip().startswith('**') and line.strip().endswith('**'):
                text = line.strip().strip('*')
                p.text = text
                p.font.size = Pt(20)
                p.font.bold = True
            else:
                p.text = line.strip()
                p.font.size = Pt(18)

    # Save presentation
    prs.save(output_file)
    print(f"PowerPoint presentation saved to: {output_file}")

def main():
    # Check for command-line argument
    if len(sys.argv) < 2:
        print("Usage: python md_to_pptx.py <input.md>")
        print("Example: python md_to_pptx.py llm-simple.md")
        sys.exit(1)

    input_file = sys.argv[1]

    # Check if file exists
    if not os.path.exists(input_file):
        print(f"Error: File '{input_file}' not found")
        sys.exit(1)

    # Generate output filename by replacing .md with .pptx
    if input_file.endswith('.md'):
        output_file = input_file[:-3] + '.pptx'
    else:
        output_file = input_file + '.pptx'

    print(f"Reading {input_file}...")
    slides = parse_markdown_slides(input_file)
    print(f"Found {len(slides)} slides")

    print(f"Creating PowerPoint presentation...")
    create_presentation(slides, output_file)
    print("Done!")

if __name__ == '__main__':
    main()
